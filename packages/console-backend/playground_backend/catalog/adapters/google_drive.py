"""Google Drive source adapter for catalog sync.

Handles file listing, page/slide extraction, thumbnail generation, and change detection
for Google Drive Shared Drives. Supports Google Slides, PPTX, PDF, and Google Docs.
"""

from __future__ import annotations

import asyncio
import io
import logging
import os
import shutil
import subprocess
import tempfile
import time
from collections.abc import Awaitable, Callable
from datetime import datetime, timezone
from typing import Any

import httplib2
import httpx
from google.oauth2.credentials import Credentials
from google_auth_httplib2 import AuthorizedHttp
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from pdf2image import convert_from_path
from pdf2image.pdf2image import pdfinfo_from_path

from ..executor import get_sync_executor, run_in_sync_executor
from .base import (
    CatalogSourceAdapter,
    ChangeSet,
    ExtractedPage,
    SourceFile,
)

logger = logging.getLogger(__name__)


class ExportTimeoutError(Exception):
    """Raised when a Drive export repeatedly times out, indicating the file is too large to export within the timeout window."""


class _SlidesRateLimiter:
    """Async rate limiter for Google Slides API (60 expensive reads/min/user).

    Serialises all Slides thumbnail requests across concurrent file syncs
    and spaces them to stay under the quota.  Uses a limit of
    58 req/min (≈ 1.03 s between calls) to leave ~3% headroom.
    """

    def __init__(self, max_per_minute: int = 58) -> None:
        self._lock = asyncio.Lock()
        self._interval = 60.0 / max_per_minute
        self._last_call = 0.0

    async def acquire(self) -> None:
        async with self._lock:
            now = time.monotonic()
            wait = self._interval - (now - self._last_call)
            if wait > 0:
                await asyncio.sleep(wait)
            self._last_call = time.monotonic()


_slides_rate_limiter = _SlidesRateLimiter()

# Shared httpx client for downloading Slides API thumbnail content URLs.
# Reused across all calls to amortise TCP/TLS handshake cost.
_thumbnail_http_client: httpx.AsyncClient | None = None


def _get_thumbnail_http_client() -> httpx.AsyncClient:
    """Return the shared httpx client, creating it lazily."""
    global _thumbnail_http_client
    if _thumbnail_http_client is None or _thumbnail_http_client.is_closed:
        _thumbnail_http_client = httpx.AsyncClient(
            timeout=httpx.Timeout(10.0, connect=5.0),
            limits=httpx.Limits(max_connections=20, max_keepalive_connections=10),
        )
    return _thumbnail_http_client


# MIME types we can process
SUPPORTED_MIME_TYPES = {
    "application/vnd.google-apps.presentation",  # Google Slides
    "application/vnd.google-apps.document",  # Google Docs
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # PPTX
    "application/pdf",
}

# Google-native MIME types (require export, not direct download)
GOOGLE_NATIVE_TYPES = {
    "application/vnd.google-apps.presentation",
    "application/vnd.google-apps.document",
}

# Binary .pptx thumbnails require spawning soffice (~700 MB subprocess RSS)
# on top of python-pptx's lxml fragmentation, which has historically pushed
# the worker into OOM under the cgroup memory limit. Off by default; set
# CATALOG_PPTX_THUMBNAILS=1 to re-enable once the worker has enough headroom.
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PPTX_THUMBNAILS_ENABLED = os.environ.get("CATALOG_PPTX_THUMBNAILS", "0") == "1"

# Fields to request from Drive API
_DRIVE_FILE_FIELDS = "id, name, mimeType, modifiedTime, parents, size, owners, webViewLink"


# Timeout for Drive API HTTP requests (seconds).
# Short timeout for metadata/list calls; longer for file downloads/exports.
_DRIVE_API_TIMEOUT = int(os.environ.get("CATALOG_DRIVE_API_TIMEOUT", "30"))
_DRIVE_EXPORT_TIMEOUT = int(os.environ.get("CATALOG_DRIVE_EXPORT_TIMEOUT", "120"))

# Maximum concurrent Drive export operations (PDF/PPTX).
# Too many parallel exports overwhelm Google's API causing mass timeouts.
_MAX_CONCURRENT_EXPORTS = int(os.environ.get("CATALOG_MAX_CONCURRENT_EXPORTS", "3"))
_export_semaphore: asyncio.Semaphore | None = None


def _get_export_semaphore() -> asyncio.Semaphore:
    """Return the shared export semaphore, creating it lazily."""
    global _export_semaphore
    if _export_semaphore is None:
        _export_semaphore = asyncio.Semaphore(_MAX_CONCURRENT_EXPORTS)
    return _export_semaphore


# Path to LibreOffice/soffice binary for PPTX→PDF conversion.
# Resolved once on first use.
_SOFFICE_PATH: str | None | bool = None  # None = not resolved yet, False = not found


def _find_soffice() -> str | None:
    """Find the LibreOffice soffice binary on this system."""
    global _SOFFICE_PATH
    if _SOFFICE_PATH is not None:
        return _SOFFICE_PATH if _SOFFICE_PATH else None

    # Check common locations
    path = shutil.which("soffice") or shutil.which("libreoffice")
    if not path:
        # macOS app bundle (system-wide or user-local)
        for mac_path in (
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ):
            if os.path.isfile(mac_path):
                path = mac_path
                break

    _SOFFICE_PATH = path or False
    if path:
        logger.info("Found LibreOffice at %s", path)
    else:
        logger.warning("LibreOffice not found — PPTX thumbnails will be unavailable")
    return path if path else None


async def _convert_office_to_pdf_to_path(input_path: str, output_dir: str) -> str:
    """Convert an Office document to PDF using LibreOffice headless.

    Reads from ``input_path`` and writes the PDF into ``output_dir``.
    Returns the path of the produced PDF. Avoids loading the file into the
    Python heap on either side of the soffice call.

    Raises RuntimeError if LibreOffice is not available.
    """
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found — cannot convert PPTX to PDF for thumbnails")

    def _convert() -> str:
        subprocess.run(
            [soffice, "--headless", "--norestore", "--convert-to", "pdf", "--outdir", output_dir, input_path],
            check=True,
            timeout=120,
            capture_output=True,
        )
        # soffice names the output after the input basename
        base = os.path.splitext(os.path.basename(input_path))[0]
        return os.path.join(output_dir, f"{base}.pdf")

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(get_sync_executor(), _convert)


def _build_drive_service(credentials: Credentials, timeout: int = _DRIVE_API_TIMEOUT) -> Any:
    """Build a Google Drive API v3 service."""
    http = AuthorizedHttp(credentials, http=httplib2.Http(timeout=timeout))
    return build("drive", "v3", http=http, cache_discovery=False)


def _build_slides_service(credentials: Credentials) -> Any:
    """Build a Google Slides API v1 service."""
    return build("slides", "v1", credentials=credentials, cache_discovery=False)


def _parse_drive_time(time_str: str) -> datetime:
    """Parse Google Drive datetime string to timezone-aware datetime."""
    return datetime.fromisoformat(time_str.replace("Z", "+00:00"))


def _is_retryable_http_error(exc: HttpError) -> bool:
    """Check if an HttpError is transient and worth retrying."""
    status = exc.resp.status if hasattr(exc, "resp") and exc.resp else 0
    return status in (429, 500, 502, 503)


async def _run_in_executor(func: Any, *args: Any) -> Any:
    """Run a blocking Google API call in a thread executor with retry for transient errors.

    Retries on ConnectionRefusedError, ConnectionResetError, OSError,
    httplib2.ServerNotFoundError (transient network/DNS issues),
    TimeoutError (read timeouts), and transient Google API HttpErrors
    (429, 500, 502, 503) with exponential backoff.

    Timeouts get only 1 retry (they indicate the file is too large to
    export within the deadline, so further retries are futile).  After
    a second timeout, ``ExportTimeoutError`` is raised so callers can
    fall back to an alternative strategy instead of retrying.
    """
    loop = asyncio.get_running_loop()
    executor = get_sync_executor()
    max_retries = 3
    timeout_count = 0
    for attempt in range(max_retries + 1):
        try:
            return await loop.run_in_executor(executor, func, *args)
        except HttpError as exc:
            if not _is_retryable_http_error(exc) or attempt == max_retries:
                raise
            wait = (2**attempt) + (0.5 * attempt)  # 1s, 2.5s, 4.5s
            logger.warning(
                "Transient Google API error %d (attempt %d/%d), retrying in %.1fs: %s",
                exc.resp.status,
                attempt + 1,
                max_retries + 1,
                wait,
                exc,
            )
            await asyncio.sleep(wait)
        except (ConnectionRefusedError, ConnectionResetError, OSError, httplib2.ServerNotFoundError) as exc:
            is_timeout = "timed out" in str(exc).lower()
            if is_timeout:
                timeout_count += 1
            if attempt == max_retries:
                if is_timeout:
                    raise ExportTimeoutError(f"Export timed out {timeout_count} times for this request") from exc
                raise
            # Timeouts only get 1 retry — a second timeout means the file
            # is too large to export within the deadline.
            if timeout_count >= 2:
                raise ExportTimeoutError(f"Export timed out {timeout_count} times, giving up") from exc
            wait = (2**attempt) + (0.5 * attempt)  # 1s, 2.5s, 4.5s
            logger.warning(
                "Transient connection error (attempt %d/%d), retrying in %.1fs: %s",
                attempt + 1,
                max_retries + 1,
                wait,
                exc,
            )
            await asyncio.sleep(wait)


def _build_folder_path(file_id: str, parents_map: dict[str, dict]) -> str:
    """Build the folder path by walking up the parent chain."""
    parts: list[str] = []
    current_id = file_id
    visited: set[str] = set()
    while current_id in parents_map and current_id not in visited:
        visited.add(current_id)
        info = parents_map[current_id]
        if info.get("parent_id"):
            parts.append(info["name"])
            current_id = info["parent_id"]
        else:
            break
    parts.reverse()
    return "/".join(parts) if parts else ""


class GoogleDriveAdapter(CatalogSourceAdapter):
    """Source adapter for Google Drive Shared Drives and shared folders."""

    async def list_files(
        self,
        config: dict[str, Any],
        progress_callback: Callable[[str], Awaitable[None]] | None = None,
    ) -> list[SourceFile]:
        """List all eligible files from a source.

        Dispatches based on source ``type``:
        - ``shared_drive`` / ``drive_folder``: query the entire Shared Drive
          (with optional folder post-filter).
        - ``shared_folder``: recursively list files under a user-shared folder.

        Legacy configs without a ``type`` key are treated as shared_drive/drive_folder.
        """
        source_type = config.get("type", "shared_drive")
        if source_type == "shared_folder":
            return await self._list_files_from_shared_folder(config, progress_callback)
        return await self._list_files_from_shared_drive(config)

    # ------------------------------------------------------------------
    # Shared Drive listing (original behaviour)
    # ------------------------------------------------------------------

    async def _list_files_from_shared_drive(self, config: dict[str, Any]) -> list[SourceFile]:
        """List eligible files from a Shared Drive, optionally filtered by folder."""
        credentials = config["credentials"]
        shared_drive_id = config["drive_id"] if "drive_id" in config else config["shared_drive_id"]
        folder_id = config.get("folder_id")
        exclude_patterns = [p.lower() for p in config.get("exclude_folder_patterns", [])]

        service = _build_drive_service(credentials)

        # Build query for supported file types
        mime_clauses = " or ".join(f"mimeType='{mt}'" for mt in SUPPORTED_MIME_TYPES)
        query = f"({mime_clauses}) and trashed=false"

        # Build a folder map for path resolution
        folder_map = await self._build_folder_map(service, shared_drive_id)

        # Pre-compute excluded folder IDs (folders whose name matches any pattern)
        excluded_folder_ids: set[str] = set()
        if exclude_patterns:
            for fid, info in folder_map.items():
                name_lower = info["name"].lower()
                if any(pat in name_lower for pat in exclude_patterns):
                    excluded_folder_ids.add(fid)

        files: list[SourceFile] = []
        page_token = None
        while True:
            request = service.files().list(
                q=query,
                corpora="drive",
                driveId=shared_drive_id,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                fields=f"nextPageToken, files({_DRIVE_FILE_FIELDS})",
                pageSize=100,
                pageToken=page_token,
            )
            result = await _run_in_executor(request.execute)

            for f in result.get("files", []):
                if folder_id:
                    parent_ids = f.get("parents", [])
                    if not self._is_under_folder(parent_ids, folder_id, folder_map):
                        continue

                # Skip files under excluded folders
                if excluded_folder_ids:
                    parent_ids = f.get("parents", [])
                    if self._is_under_any_folder(parent_ids, excluded_folder_ids, folder_map):
                        continue

                parent_id = f.get("parents", [None])[0]
                folder_path = _build_folder_path(parent_id, folder_map) if parent_id else ""

                files.append(self._make_source_file(f, folder_path))

            page_token = result.get("nextPageToken")
            if not page_token:
                break

        logger.info("Listed %d files from Shared Drive %s", len(files), shared_drive_id)
        return files

    # ------------------------------------------------------------------
    # User-shared folder listing — two-phase BFS approach
    # ------------------------------------------------------------------

    async def _list_files_from_shared_folder(
        self,
        config: dict[str, Any],
        progress_callback: Callable[[str], Awaitable[None]] | None = None,
    ) -> list[SourceFile]:
        """List eligible files under a user-shared folder.

        Uses a two-phase approach for efficient traversal:
          **Phase 1 — Folder discovery**: BFS with parallel queries per level
          to discover all subfolder IDs and build a ``{folder_id: path}`` map.
          Folder-only queries are lightweight (``pageSize=1000``).
          **Phase 2 — File listing**: Query all discovered folders for files
          in parallel.  Since the full folder tree is known, every folder
          can be queried simultaneously (bounded by a concurrency semaphore).

        A fixed-size pool of Drive service objects is shared across tasks
        so that at most ``_DRIVE_POOL_SIZE`` ``httplib2.Http`` instances
        exist at any time, preventing OOM on large folder trees.
        """
        credentials = config["credentials"]
        folder_id = config["folder_id"]
        folder_name = config.get("folder_name", "")
        exclude_patterns = [p.lower() for p in config.get("exclude_folder_patterns", [])]

        # Pool of reusable Drive service objects.  Each httplib2.Http is NOT
        # thread-safe, so we let at most _DRIVE_POOL_SIZE tasks use the API
        # concurrently — one service per task while it holds the slot.
        _DRIVE_POOL_SIZE = 10
        svc_pool: asyncio.Queue[Any] = asyncio.Queue(maxsize=_DRIVE_POOL_SIZE)
        for _ in range(_DRIVE_POOL_SIZE):
            svc_pool.put_nowait(_build_drive_service(credentials))

        try:
            # -- Phase 1: BFS folder discovery (parallel per level) --------
            # folder_map maps every folder_id → its human-readable path
            folder_map: dict[str, str] = {folder_id: folder_name}
            current_level = [folder_id]

            while current_level:

                async def _list_child_folders(parent_id: str) -> list[dict]:
                    """Return child folder dicts for one parent."""
                    children: list[dict] = []
                    page_token = None
                    svc = await svc_pool.get()
                    try:
                        while True:
                            request = svc.files().list(
                                q=(
                                    f"mimeType='application/vnd.google-apps.folder' "
                                    f"and '{parent_id}' in parents and trashed=false"
                                ),
                                includeItemsFromAllDrives=True,
                                supportsAllDrives=True,
                                fields="nextPageToken, files(id, name)",
                                pageSize=1000,
                                pageToken=page_token,
                            )
                            result = await _run_in_executor(request.execute)
                            children.extend(result.get("files", []))
                            page_token = result.get("nextPageToken")
                            if not page_token:
                                break
                    finally:
                        svc_pool.put_nowait(svc)
                    return children

                # Query all parents in this level in parallel (bounded by pool size)
                level_results = await asyncio.gather(*[_list_child_folders(pid) for pid in current_level])

                next_level: list[str] = []
                for parent_id, children in zip(current_level, level_results):
                    parent_path = folder_map[parent_id]
                    for child in children:
                        name_lower = child["name"].lower()
                        if exclude_patterns and any(pat in name_lower for pat in exclude_patterns):
                            logger.debug("Excluding folder %s (matched exclusion pattern)", child["name"])
                            continue
                        child_path = f"{parent_path}/{child['name']}" if parent_path else child["name"]
                        folder_map[child["id"]] = child_path
                        next_level.append(child["id"])
                current_level = next_level

                if progress_callback:
                    await progress_callback(f"Discovering folders... ({len(folder_map):,} found)")

            logger.info(
                "Discovered %d folders under %s (%s)",
                len(folder_map),
                folder_name,
                folder_id,
            )

            # -- Phase 2: List files from ALL folders in parallel ----------
            mime_clauses = " or ".join(f"mimeType='{mt}'" for mt in SUPPORTED_MIME_TYPES)
            file_count = 0  # shared counter — safe because only one asyncio task runs at a time

            async def _list_files_in_folder(fid: str) -> list[SourceFile]:
                """List eligible (non-folder) files in a single folder."""
                nonlocal file_count
                fpath = folder_map[fid]
                found: list[SourceFile] = []
                page_token = None
                svc = await svc_pool.get()
                try:
                    while True:
                        request = svc.files().list(
                            q=f"({mime_clauses}) and '{fid}' in parents and trashed=false",
                            includeItemsFromAllDrives=True,
                            supportsAllDrives=True,
                            fields=f"nextPageToken, files({_DRIVE_FILE_FIELDS})",
                            pageSize=100,
                            pageToken=page_token,
                        )
                        result = await _run_in_executor(request.execute)
                        for f in result.get("files", []):
                            found.append(self._make_source_file(f, fpath))
                        page_token = result.get("nextPageToken")
                        if not page_token:
                            break
                finally:
                    svc_pool.put_nowait(svc)
                file_count += len(found)
                if progress_callback:
                    await progress_callback(f"Listing files... ({file_count:,} found)")
                return found

            file_results = await asyncio.gather(*[_list_files_in_folder(fid) for fid in folder_map])
            all_files = [f for batch in file_results for f in batch]

            logger.info(
                "Listed %d files from shared folder %s (%s)",
                len(all_files),
                folder_name,
                folder_id,
            )
            return all_files
        finally:
            # Close all pooled service objects to release HTTP connections.
            while not svc_pool.empty():
                svc_pool.get_nowait().close()

    @staticmethod
    def _make_source_file(f: dict, folder_path: str) -> SourceFile:
        """Create a SourceFile from a Drive API file dict."""
        return SourceFile(
            id=f["id"],
            name=f["name"],
            mime_type=f["mimeType"],
            modified_at=_parse_drive_time(f["modifiedTime"]),
            folder_path=folder_path,
            metadata={
                "size": f.get("size"),
                "owners": [o.get("displayName", "") for o in f.get("owners", [])],
                "web_view_link": f.get("webViewLink", ""),
            },
        )

    async def extract_pages(self, file: SourceFile, credentials: Any) -> list[ExtractedPage]:
        """Extract pages from a file based on its MIME type."""
        if file.mime_type == "application/vnd.google-apps.presentation":
            return await self._extract_google_slides(file, credentials)
        elif file.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return await self._extract_pptx(file, credentials)
        elif file.mime_type == "application/pdf":
            return await self._extract_pdf(file, credentials)
        elif file.mime_type == "application/vnd.google-apps.document":
            return await self._extract_google_doc(file, credentials)
        else:
            logger.warning("Unsupported MIME type %s for file %s", file.mime_type, file.name)
            return []

    async def get_thumbnail(self, file: SourceFile, page: ExtractedPage, credentials: Any) -> bytes | None:
        """Get thumbnail for a page/slide.

        Binary .pptx thumbnails are gated behind ``CATALOG_PPTX_THUMBNAILS=1``
        because rendering them requires spawning soffice (~700 MB subprocess
        RSS) on top of python-pptx's lxml fragmentation, which can push the
        worker into OOM under tight cgroup memory limits. When disabled,
        text and summary extraction still run; users who want thumbnails
        should upload as Google Slides.
        """
        if file.mime_type == "application/vnd.google-apps.presentation":
            return await self._get_slides_thumbnail(file, page, credentials)
        elif file.mime_type in (
            "application/pdf",
            "application/vnd.google-apps.document",
        ):
            return await self._get_pdf_page_thumbnail(file, page, credentials)
        elif file.mime_type == PPTX_MIME and _PPTX_THUMBNAILS_ENABLED:
            return await self._get_pdf_page_thumbnail(file, page, credentials)
        return None

    async def get_all_thumbnails(
        self,
        file: SourceFile,
        pages: list[ExtractedPage],
        credentials: Any,
    ) -> dict[int, bytes]:
        """Get thumbnails for all pages, downloading PDF-based files only once.

        Binary .pptx is gated behind ``CATALOG_PPTX_THUMBNAILS=1``; see
        :meth:`get_thumbnail` for the rationale.
        """
        thumbnail_types = {
            "application/vnd.google-apps.presentation",
            "application/pdf",
            "application/vnd.google-apps.document",
        }
        if _PPTX_THUMBNAILS_ENABLED:
            thumbnail_types.add(PPTX_MIME)
        if file.mime_type in thumbnail_types:
            return await self._get_all_pdf_thumbnails(file, pages, credentials)

        return {}

    async def _get_all_pdf_thumbnails(
        self,
        file: SourceFile,
        pages: list[ExtractedPage],
        credentials: Any,
    ) -> dict[int, bytes]:
        """Download file once, render all page thumbnails from disk.

        Streams the file directly to disk (never buffered in Python heap) and
        runs soffice / poppler on the on-disk path. For Google Slides that
        exceed the 10 MB PDF export limit, falls back to per-slide Slides API
        thumbnails (rate-limited but avoids the cap).
        """
        # Single temp directory for the file's whole render lifecycle.
        # All downloads / soffice output / poppler input stay on disk inside it.
        tmpdir = tempfile.mkdtemp(prefix="catalog-render-")
        try:
            pdf_path: str | None = None

            if file.mime_type in GOOGLE_NATIVE_TYPES:
                pdf_path = os.path.join(tmpdir, "input.pdf")
                try:
                    await self._export_to_path(file.id, "application/pdf", credentials, pdf_path)
                except (HttpError, ExportTimeoutError) as exc:
                    is_size_limit = isinstance(exc, HttpError) and "exportSizeLimitExceeded" in str(exc)
                    is_timeout = isinstance(exc, ExportTimeoutError)
                    if file.mime_type == "application/vnd.google-apps.presentation" and (is_size_limit or is_timeout):
                        # Drive's files.export enforces a 10 MB cap on Slides
                        # PDF exports. Try the exportLinks URL first — same
                        # endpoint the Drive web UI uses, no size cap, streams
                        # straight to disk so peak memory stays bounded.
                        # Only fall back to the per-slide Slides API path
                        # (rate-limited at 58/min) if exportLinks also fails.
                        reason = "timeout" if is_timeout else "size limit"
                        try:
                            logger.info(
                                "File %s PDF export failed (%s), retrying via exportLinks",
                                file.name,
                                reason,
                            )
                            await self._export_via_link_to_path(file.id, credentials, pdf_path)
                        except Exception as link_exc:
                            logger.info(
                                "File %s exportLinks fallback failed (%s), using Slides API thumbnails",
                                file.name,
                                link_exc,
                            )
                            return await self._get_all_slides_thumbnails(file, pages, credentials)
                    elif is_timeout:
                        # Non-Slides Google-native file (e.g. Docs) that timed out:
                        # skip thumbnails entirely rather than retrying forever.
                        logger.warning(
                            "File %s export timed out, skipping thumbnails",
                            file.name,
                        )
                        return {}
                    else:
                        raise
            elif file.mime_type == PPTX_MIME:
                # Reached only when CATALOG_PPTX_THUMBNAILS=1; otherwise the
                # dispatch in get_all_thumbnails() short-circuits with {}.
                pptx_path = os.path.join(tmpdir, "input.pptx")
                await self._download_file_to_path(file.id, credentials, pptx_path)
                pdf_path = await _convert_office_to_pdf_to_path(pptx_path, tmpdir)
                os.unlink(pptx_path)
            else:
                pdf_path = os.path.join(tmpdir, "input.pdf")
                await self._download_file_to_path(file.id, credentials, pdf_path)

            assert pdf_path is not None

            # Render pages in batches to bound peak memory.
            # Each batch spawns poppler for a page range, converts to PIL,
            # extracts needed PNGs, then frees the PIL objects before next batch.
            _BATCH_SIZE = 10
            # 72 DPI ≈ 960×540 px for 16:9 slides — sufficient for 800px-wide thumbnails.
            # PIL holds RGB bitmaps in memory: 960×540×3 = ~1.5 MB per page × batch=10 = ~15 MB peak.
            # Doubling DPI quadruples PIL memory, so 72 is the sweet spot for our use case.
            _THUMB_DPI = 72

            result: dict[int, bytes] = {}
            needed = {p.page_number for p in pages}

            info = await run_in_sync_executor(
                lambda: pdfinfo_from_path(pdf_path),
            )
            total_pages = info.get("Pages", 0)
            if total_pages == 0:
                return result

            # Process all pages in batches
            for batch_start in range(1, total_pages + 1, _BATCH_SIZE):
                batch_end = min(batch_start + _BATCH_SIZE - 1, total_pages)

                # Check if any needed pages are in this range
                batch_needed = needed & set(range(batch_start, batch_end + 1))
                if not batch_needed:
                    continue

                _s = batch_start
                _e = batch_end
                _path = pdf_path
                images = await run_in_sync_executor(
                    lambda: convert_from_path(
                        _path,
                        dpi=_THUMB_DPI,
                        fmt="png",
                        first_page=_s,
                        last_page=_e,
                    ),
                )

                for i, img in enumerate(images):
                    page_num = batch_start + i
                    if page_num in needed:
                        buf = io.BytesIO()
                        img.save(buf, format="PNG")
                        result[page_num] = buf.getvalue()
                    img.close()

        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

        return result

    async def _get_all_slides_thumbnails(
        self,
        file: SourceFile,
        pages: list[ExtractedPage],
        credentials: Any,
    ) -> dict[int, bytes]:
        """Get thumbnails for all slides via the Slides API (per-slide).

        Fallback for presentations that exceed the Drive PDF export size limit
        AND the PPTX export fallback.  Each API call goes through
        ``_slides_rate_limiter`` to stay within quota.  Fetch tasks run
        concurrently so that the httpx content-URL download for one slide
        overlaps with the rate-limiter wait for the next.
        """
        result: dict[int, bytes] = {}

        async def _fetch_one(page: ExtractedPage) -> tuple[int, bytes | None]:
            try:
                thumb = await self._get_slides_thumbnail(file, page, credentials)
                return page.page_number, thumb
            except Exception:
                logger.warning(
                    "Slides API thumbnail failed for page %d of %s",
                    page.page_number,
                    file.name,
                    exc_info=True,
                )
                return page.page_number, None

        fetched = await asyncio.gather(*[_fetch_one(page) for page in pages])
        for page_number, thumb in fetched:
            if thumb:
                result[page_number] = thumb

        return result

    async def detect_changes(self, config: dict[str, Any], since_token: str | None = None) -> ChangeSet:
        """Detect changes using Drive API changes.list().

        Dispatches based on source ``type``:
        - ``shared_drive`` / ``drive_folder``: drive-scoped changes API.
        - ``shared_folder``: user-scoped changes API with folder ancestry filter.
        """
        source_type = config.get("type", "shared_drive")
        credentials = config["credentials"]
        service = _build_drive_service(credentials)

        is_drive_source = source_type in ("shared_drive", "drive_folder")
        drive_id = config.get("drive_id") or config.get("shared_drive_id") if is_drive_source else None
        folder_id = config.get("folder_id")

        # ----- get initial token -----
        if since_token is None:
            if drive_id:
                request = service.changes().getStartPageToken(driveId=drive_id, supportsAllDrives=True)
            else:
                request = service.changes().getStartPageToken(supportsAllDrives=True)
            result = await _run_in_executor(request.execute)
            return ChangeSet(new_page_token=result["startPageToken"])

        # ----- build optional folder map for ancestry check -----
        folder_map: dict[str, dict] | None = None
        if folder_id and drive_id:
            folder_map = await self._build_folder_map(service, drive_id)

        added: list[SourceFile] = []
        modified: list[SourceFile] = []
        deleted_ids: list[str] = []
        page_token = since_token

        while page_token:
            list_kwargs: dict[str, Any] = {
                "pageToken": page_token,
                "includeItemsFromAllDrives": True,
                "supportsAllDrives": True,
                "fields": "nextPageToken, newStartPageToken, changes(fileId, removed, file(id, name, mimeType, modifiedTime, parents, trashed))",
                "pageSize": 100,
            }
            if drive_id:
                list_kwargs["driveId"] = drive_id

            request = service.changes().list(**list_kwargs)
            result = await _run_in_executor(request.execute)

            for change in result.get("changes", []):
                file_id = change["fileId"]
                if change.get("removed") or (change.get("file", {}).get("trashed")):
                    deleted_ids.append(file_id)
                    continue

                f = change.get("file", {})
                mime_type = f.get("mimeType", "")
                if mime_type not in SUPPORTED_MIME_TYPES:
                    continue

                # Apply folder filter when source is a subfolder
                if folder_id:
                    parents = f.get("parents", [])
                    if folder_map:
                        if not self._is_under_folder(parents, folder_id, folder_map):
                            continue
                    elif not self._is_direct_or_recursive_child(parents, folder_id):
                        continue

                source_file = SourceFile(
                    id=f["id"],
                    name=f.get("name", ""),
                    mime_type=mime_type,
                    modified_at=_parse_drive_time(f["modifiedTime"])
                    if f.get("modifiedTime")
                    else datetime.now(timezone.utc),
                )
                modified.append(source_file)

            page_token = result.get("nextPageToken")
            if not page_token:
                new_token = result.get("newStartPageToken")
                return ChangeSet(
                    added=added,
                    modified=modified,
                    deleted_ids=deleted_ids,
                    new_page_token=new_token,
                )

        return ChangeSet(added=added, modified=modified, deleted_ids=deleted_ids)

    @staticmethod
    def _is_direct_or_recursive_child(parent_ids: list[str], target_folder_id: str) -> bool:
        """Simple check: is the file a direct child of the target?

        For user-shared folders we don't have a full folder map, so we fall back
        to checking only the first parent.  The recursive listing already
        captured all descendant files during ``list_files``, so incremental
        deletes outside the folder are harmless (file won't exist in DB).
        """
        return target_folder_id in parent_ids

    async def list_shared_drives(self, credentials: Any) -> list[dict[str, str]]:
        """List shared drives accessible to the user."""
        service = _build_drive_service(credentials)
        drives: list[dict[str, str]] = []
        page_token = None

        while True:
            request = service.drives().list(
                pageSize=100,
                pageToken=page_token,
                fields="nextPageToken, drives(id, name)",
            )
            result = await _run_in_executor(request.execute)
            for d in result.get("drives", []):
                drives.append({"id": d["id"], "name": d["name"]})

            page_token = result.get("nextPageToken")
            if not page_token:
                break

        return drives

    async def list_folders(
        self,
        credentials: Any,
        shared_drive_id: str | None = None,
        parent_id: str | None = None,
    ) -> list[dict[str, str]]:
        """List direct child folders.

        When *shared_drive_id* is given the query is scoped to that Shared
        Drive (``corpora="drive"``).  Otherwise, a plain ``files.list``
        call is used which works for any folder the caller has access to
        (e.g. folders from 'Shared with me').

        Args:
            credentials: Google OAuth credentials.
            shared_drive_id: Optional Shared Drive ID.
            parent_id: Parent folder ID.  Required when *shared_drive_id*
                       is ``None``.  When *shared_drive_id* is given and
                       *parent_id* is ``None`` the root of the drive is
                       listed.

        Returns:
            List of dicts with ``id`` and ``name`` keys, sorted by name.
        """
        service = _build_drive_service(credentials)

        if shared_drive_id:
            parent = parent_id or shared_drive_id
        else:
            if not parent_id:
                return []  # cannot list root without a drive
            parent = parent_id

        query = f"mimeType='application/vnd.google-apps.folder' and '{parent}' in parents and trashed=false"

        extra_kwargs: dict[str, Any] = {}
        if shared_drive_id:
            extra_kwargs.update(corpora="drive", driveId=shared_drive_id)

        folders: list[dict[str, str]] = []
        page_token = None

        while True:
            request = service.files().list(
                q=query,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                fields="nextPageToken, files(id, name)",
                pageSize=1000,
                orderBy="name",
                pageToken=page_token,
                **extra_kwargs,
            )
            result = await _run_in_executor(request.execute)
            for f in result.get("files", []):
                folders.append({"id": f["id"], "name": f["name"]})

            page_token = result.get("nextPageToken")
            if not page_token:
                break

        return folders

    async def list_user_shared_folders(self, credentials: Any) -> list[dict[str, str]]:
        """List folders shared with the user (from 'Shared with me').

        Returns:
            List of dicts with ``id`` and ``name`` keys, sorted by name.
        """
        service = _build_drive_service(credentials)
        query = "sharedWithMe=true and mimeType='application/vnd.google-apps.folder' and trashed=false"
        folders: list[dict[str, str]] = []
        page_token = None

        while True:
            request = service.files().list(
                q=query,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                fields="nextPageToken, files(id, name)",
                pageSize=100,
                orderBy="name",
                pageToken=page_token,
            )
            result = await _run_in_executor(request.execute)
            for f in result.get("files", []):
                folders.append({"id": f["id"], "name": f["name"]})

            page_token = result.get("nextPageToken")
            if not page_token:
                break

        return folders

    # --- Private extraction methods ---

    async def _extract_google_slides(self, file: SourceFile, credentials: Credentials) -> list[ExtractedPage]:
        """Extract pages from Google Slides using the Slides API."""
        service = _build_slides_service(credentials)
        request = service.presentations().get(presentationId=file.id)
        presentation = await _run_in_executor(request.execute)

        pages: list[ExtractedPage] = []
        for idx, slide in enumerate(presentation.get("slides", []), start=1):
            text_parts: list[str] = []
            title = ""
            notes = ""

            # Extract text from page elements
            for element in slide.get("pageElements", []):
                shape = element.get("shape", {})
                text_obj = shape.get("text", {})
                placeholder = shape.get("placeholder", {})

                element_text = self._extract_text_from_text_elements(text_obj)

                if placeholder.get("type") == "TITLE":
                    title = element_text.strip()
                elif element_text.strip():
                    text_parts.append(element_text.strip())

            # Extract speaker notes
            notes_page = slide.get("slideProperties", {}).get("notesPage", {})
            for element in notes_page.get("pageElements", []):
                shape = element.get("shape", {})
                placeholder = shape.get("placeholder", {})
                if placeholder.get("type") == "BODY":
                    text_obj = shape.get("text", {})
                    notes = self._extract_text_from_text_elements(text_obj).strip()

            pages.append(
                ExtractedPage(
                    page_number=idx,
                    title=title or f"Slide {idx}",
                    text_content="\n".join(text_parts),
                    speaker_notes=notes,
                    source_ref={
                        "type": "google_slides",
                        "file_id": file.id,
                        "page_object_id": slide.get("objectId", ""),
                        "presentation_id": file.id,
                    },
                )
            )

        return pages

    async def _extract_pptx(self, file: SourceFile, credentials: Credentials) -> list[ExtractedPage]:
        """Extract pages from a PPTX file using python-pptx.

        Streams the download directly to a temp file and lets python-pptx open
        the path. Avoids holding the full file as a Python ``bytes`` plus an
        ``io.BytesIO`` copy in heap (which fragments the lxml allocator and
        retains 400+ MB of RSS even after parsing completes).
        """
        from pptx import Presentation

        tmpdir = tempfile.mkdtemp(prefix="catalog-pptx-")
        pptx_path = os.path.join(tmpdir, "input.pptx")
        try:
            await self._download_file_to_path(file.id, credentials, pptx_path)

            def _parse() -> list[ExtractedPage]:
                prs = Presentation(pptx_path)

                pages: list[ExtractedPage] = []
                visible_idx = 0
                for original_idx, slide in enumerate(prs.slides):
                    # Skip hidden slides — LibreOffice omits them from the PDF,
                    # so page numbering must match to keep thumbnails aligned.
                    if slide._element.get("show") == "0":
                        continue
                    visible_idx += 1

                    text_parts: list[str] = []
                    title = ""

                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if (
                                shape.shape_id == slide.shapes.title
                                and hasattr(slide.shapes, "title")
                                and slide.shapes.title == shape
                            ):
                                title = text
                            elif text:
                                text_parts.append(text)

                    # Try to get title from placeholders
                    if not title and slide.shapes.title:
                        title = slide.shapes.title.text.strip()

                    # Speaker notes
                    notes = ""
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()

                    pages.append(
                        ExtractedPage(
                            page_number=visible_idx,
                            title=title or f"Slide {visible_idx}",
                            text_content="\n".join(text_parts),
                            speaker_notes=notes,
                            source_ref={
                                "type": "pptx",
                                "file_id": file.id,
                                "slide_index": original_idx,  # 0-based original index (including hidden slides)
                            },
                        )
                    )

                return pages

            return await run_in_sync_executor(_parse)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    async def _extract_pdf(self, file: SourceFile, credentials: Credentials) -> list[ExtractedPage]:
        """Extract pages from a PDF file.

        Streams the download directly to a temp file and lets pypdf open the
        path. Avoids holding the full file as a Python ``bytes`` plus an
        ``io.BytesIO`` copy in heap.
        """
        import pypdf

        tmpdir = tempfile.mkdtemp(prefix="catalog-pdf-")
        pdf_path = os.path.join(tmpdir, "input.pdf")
        try:
            await self._download_file_to_path(file.id, credentials, pdf_path)

            def _parse() -> list[ExtractedPage]:
                reader = pypdf.PdfReader(pdf_path)

                pages: list[ExtractedPage] = []
                for idx, pdf_page in enumerate(reader.pages, start=1):
                    text = pdf_page.extract_text() or ""
                    pages.append(
                        ExtractedPage(
                            page_number=idx,
                            title=f"Page {idx}",
                            text_content=text.strip(),
                            source_ref={
                                "type": "pdf",
                                "file_id": file.id,
                                "page_number": idx,
                            },
                        )
                    )

                return pages

            return await run_in_sync_executor(_parse)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    async def _extract_google_doc(self, file: SourceFile, credentials: Credentials) -> list[ExtractedPage]:
        """Extract pages from Google Docs by exporting as PDF.

        Streams the PDF export directly to a temp file and lets pypdf open the
        path. Avoids holding the full PDF as a Python ``bytes`` plus an
        ``io.BytesIO`` copy in heap.
        """
        import pypdf

        tmpdir = tempfile.mkdtemp(prefix="catalog-gdoc-")
        pdf_path = os.path.join(tmpdir, "input.pdf")
        try:
            await self._export_to_path(file.id, "application/pdf", credentials, pdf_path)

            def _parse() -> list[ExtractedPage]:
                reader = pypdf.PdfReader(pdf_path)

                pages: list[ExtractedPage] = []
                for idx, pdf_page in enumerate(reader.pages, start=1):
                    text = pdf_page.extract_text() or ""
                    pages.append(
                        ExtractedPage(
                            page_number=idx,
                            title=f"Page {idx}",
                            text_content=text.strip(),
                            source_ref={
                                "type": "google_docs_pdf",
                                "file_id": file.id,
                                "page_number": idx,
                            },
                        )
                    )

                return pages

            return await run_in_sync_executor(_parse)
        finally:
            shutil.rmtree(tmpdir, ignore_errors=True)

    async def _get_slides_thumbnail(
        self, file: SourceFile, page: ExtractedPage, credentials: Credentials
    ) -> bytes | None:
        """Get thumbnail from Google Slides API."""
        page_object_id = page.source_ref.get("page_object_id")
        if not page_object_id:
            return None

        await _slides_rate_limiter.acquire()

        service = _build_slides_service(credentials)
        request = (
            service.presentations()
            .pages()
            .getThumbnail(
                presentationId=file.id,
                pageObjectId=page_object_id,
                thumbnailProperties_thumbnailSize="MEDIUM",
            )
        )
        result = await _run_in_executor(request.execute)
        content_url = result.get("contentUrl")
        if not content_url:
            return None

        client = _get_thumbnail_http_client()
        resp = await client.get(content_url)
        resp.raise_for_status()
        return resp.content

    async def _get_pdf_page_thumbnail(
        self, file: SourceFile, page: ExtractedPage, credentials: Credentials
    ) -> bytes | None:
        """Get thumbnail for a PDF page (or file exported as PDF) using pdf2image."""
        from pdf2image import convert_from_bytes

        if file.mime_type in GOOGLE_NATIVE_TYPES:
            pdf_bytes = await self._export_as_pdf(file.id, credentials)
        else:
            pdf_bytes = await self._download_file(file.id, credentials)

        images = await run_in_sync_executor(
            lambda: convert_from_bytes(
                pdf_bytes,
                dpi=150,
                fmt="png",
                first_page=page.page_number,
                last_page=page.page_number,
            ),
        )

        if not images:
            return None

        buf = io.BytesIO()
        images[0].save(buf, format="PNG")
        return buf.getvalue()

    # --- Helpers ---

    async def _download_file(self, file_id: str, credentials: Credentials) -> bytes:
        """Download a binary file from Google Drive."""
        service = _build_drive_service(credentials, timeout=_DRIVE_EXPORT_TIMEOUT)
        request = service.files().get_media(fileId=file_id, supportsAllDrives=True)
        buf = io.BytesIO()

        def _download() -> bytes:
            downloader = MediaIoBaseDownload(buf, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            return buf.getvalue()

        return await _run_in_executor(_download)

    async def _download_file_to_path(self, file_id: str, credentials: Credentials, dest_path: str) -> None:
        """Stream a binary file from Google Drive directly to disk.

        Avoids buffering the entire file in Python heap (BytesIO). Each
        ``next_chunk()`` write goes straight to the file descriptor.
        """
        service = _build_drive_service(credentials, timeout=_DRIVE_EXPORT_TIMEOUT)
        request = service.files().get_media(fileId=file_id, supportsAllDrives=True)

        def _download() -> None:
            with open(dest_path, "wb") as fh:
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()

        await _run_in_executor(_download)

    async def _export_as_pdf(self, file_id: str, credentials: Credentials) -> bytes:
        """Export a Google-native file as PDF."""
        async with _get_export_semaphore():
            service = _build_drive_service(credentials, timeout=_DRIVE_EXPORT_TIMEOUT)
            request = service.files().export_media(
                fileId=file_id,
                mimeType="application/pdf",
            )
            buf = io.BytesIO()

            def _download() -> bytes:
                downloader = MediaIoBaseDownload(buf, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()
                return buf.getvalue()

            return await _run_in_executor(_download)

    async def _export_to_path(
        self,
        file_id: str,
        mime_type: str,
        credentials: Credentials,
        dest_path: str,
    ) -> None:
        """Stream an export of a Google-native file directly to disk.

        Avoids buffering the entire export in Python heap.
        """
        async with _get_export_semaphore():
            service = _build_drive_service(credentials, timeout=_DRIVE_EXPORT_TIMEOUT)
            request = service.files().export_media(fileId=file_id, mimeType=mime_type)

            def _download() -> None:
                with open(dest_path, "wb") as fh:
                    downloader = MediaIoBaseDownload(fh, request)
                    done = False
                    while not done:
                        _, done = downloader.next_chunk()

            await _run_in_executor(_download)

    async def _export_via_link_to_path(
        self,
        file_id: str,
        credentials: Credentials,
        dest_path: str,
    ) -> None:
        """Export a Google-native file as PDF via the ``exportLinks`` URL.

        ``files.export`` enforces a 10 MB cap on Google Slides PDF exports.
        ``exportLinks`` returns the same download URL the Drive web UI uses
        for "Download as PDF", which has no such cap. We stream the response
        straight to disk to keep the file out of the Python heap.

        The Drive metadata fetch is performed via the authenticated client,
        which transparently refreshes ``credentials.token`` if expired —
        the freshly-minted bearer token is then reused for the raw HTTPS
        download below.
        """
        async with _get_export_semaphore():
            service = _build_drive_service(credentials, timeout=_DRIVE_EXPORT_TIMEOUT)
            meta = await _run_in_executor(
                service.files().get(fileId=file_id, fields="exportLinks", supportsAllDrives=True).execute
            )
            export_links = meta.get("exportLinks") or {}
            pdf_url = export_links.get("application/pdf")
            if not pdf_url:
                raise RuntimeError(f"File {file_id} has no application/pdf exportLinks URL")

            headers = {"Authorization": f"Bearer {credentials.token}"}
            timeout = httpx.Timeout(float(_DRIVE_EXPORT_TIMEOUT), connect=10.0)

            async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
                async with client.stream("GET", pdf_url, headers=headers) as resp:
                    resp.raise_for_status()
                    ctype = resp.headers.get("content-type", "").lower()
                    # The unrestricted endpoint occasionally returns an HTML
                    # interstitial (e.g. virus-scan warning, redirect page).
                    # Reject anything that isn't a PDF so the caller can fall
                    # back to the per-slide Slides API path.
                    if "pdf" not in ctype:
                        raise RuntimeError(f"exportLinks URL returned non-PDF content-type: {ctype}")
                    with open(dest_path, "wb") as fh:
                        async for chunk in resp.aiter_bytes(chunk_size=64 * 1024):
                            fh.write(chunk)

    async def _build_folder_map(self, service: Any, shared_drive_id: str) -> dict[str, dict]:
        """Build a map of folder_id -> {name, parent_id} for path resolution."""
        folder_map: dict[str, dict] = {}
        page_token = None

        while True:
            request = service.files().list(
                q="mimeType='application/vnd.google-apps.folder' and trashed=false",
                corpora="drive",
                driveId=shared_drive_id,
                includeItemsFromAllDrives=True,
                supportsAllDrives=True,
                fields="nextPageToken, files(id, name, parents)",
                pageSize=1000,
                pageToken=page_token,
            )
            result = await _run_in_executor(request.execute)

            for f in result.get("files", []):
                folder_map[f["id"]] = {
                    "name": f["name"],
                    "parent_id": f.get("parents", [None])[0],
                }

            page_token = result.get("nextPageToken")
            if not page_token:
                break

        return folder_map

    def _is_under_folder(
        self,
        parent_ids: list[str],
        target_folder_id: str,
        folder_map: dict[str, dict],
    ) -> bool:
        """Check if a file is under (descendant of) a target folder."""
        visited: set[str] = set()
        for pid in parent_ids:
            current = pid
            while current and current not in visited:
                if current == target_folder_id:
                    return True
                visited.add(current)
                parent_info = folder_map.get(current, {})
                current = parent_info.get("parent_id")
        return False

    def _is_under_any_folder(
        self,
        parent_ids: list[str],
        target_folder_ids: set[str],
        folder_map: dict[str, dict],
    ) -> bool:
        """Check if a file is under (descendant of) any of the target folders."""
        visited: set[str] = set()
        for pid in parent_ids:
            current = pid
            while current and current not in visited:
                if current in target_folder_ids:
                    return True
                visited.add(current)
                parent_info = folder_map.get(current, {})
                current = parent_info.get("parent_id")
        return False

    @staticmethod
    def _extract_text_from_text_elements(text_obj: dict) -> str:
        """Extract text from Google Slides text elements."""
        parts: list[str] = []
        for element in text_obj.get("textElements", []):
            text_run = element.get("textRun", {})
            content = text_run.get("content", "")
            if content:
                parts.append(content)
        return "".join(parts)
